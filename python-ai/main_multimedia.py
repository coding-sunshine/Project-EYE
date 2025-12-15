"""
Multi-Media FastAPI service with support for images, videos, documents, and audio.
Includes CLIP embeddings, face detection, video analysis, OCR, and audio transcription.
"""

from fastapi import FastAPI, HTTPException
from pydantic import BaseModel
from transformers import BlipProcessor, BlipForConditionalGeneration
from transformers import CLIPProcessor, CLIPModel
from transformers import AutoProcessor, AutoModelForCausalLM
import torch
from PIL import Image
import numpy as np
from pathlib import Path
import logging
import face_recognition
import cv2
from typing import List, Dict, Optional, Any
import json
from concurrent.futures import ThreadPoolExecutor, as_completed
import os

# Register HEIF/HEIC support
try:
    from pillow_heif import register_heif_opener
    register_heif_opener()
    logging.getLogger(__name__).info("HEIF/HEIC support registered")
except ImportError:
    logging.getLogger(__name__).warning("pillow-heif not available - HEIC files won't be supported")

# Try to import optional dependencies
try:
    import ollama
    import os
    OLLAMA_HOST = os.getenv('OLLAMA_HOST', 'http://ollama:11434')
    ollama_client = ollama.Client(host=OLLAMA_HOST)
    OLLAMA_AVAILABLE = True
    logging.info(f"Ollama is available at {OLLAMA_HOST}")
except ImportError:
    ollama_client = None
    OLLAMA_AVAILABLE = False
    logging.warning("Ollama not available, will use BLIP only")

# Import comprehensive image analyzer
try:
    from comprehensive_analyzer import ComprehensiveImageAnalyzer, create_analyzer
    COMPREHENSIVE_ANALYZER_AVAILABLE = True
    logging.info("ComprehensiveImageAnalyzer loaded successfully")
except ImportError as e:
    COMPREHENSIVE_ANALYZER_AVAILABLE = False
    logging.warning(f"ComprehensiveImageAnalyzer not available: {e}")

try:
    import whisper
    WHISPER_AVAILABLE = True
except ImportError:
    WHISPER_AVAILABLE = False
    logging.warning("Whisper not available, audio transcription disabled")

try:
    import pytesseract
    TESSERACT_AVAILABLE = True
except ImportError:
    TESSERACT_AVAILABLE = False
    logging.warning("Tesseract not available")

# PaddleOCR - Better accuracy, especially for complex layouts
try:
    from paddleocr import PaddleOCR
    PADDLE_OCR = None  # Lazy initialization
    PADDLEOCR_AVAILABLE = True
    logging.info("PaddleOCR is available")
except ImportError:
    PADDLE_OCR = None
    PADDLEOCR_AVAILABLE = False
    logging.warning("PaddleOCR not available")

# At least one OCR engine must be available
OCR_AVAILABLE = TESSERACT_AVAILABLE or PADDLEOCR_AVAILABLE
if not OCR_AVAILABLE:
    logging.warning("No OCR engine available, OCR disabled")

# Parallel processing configuration
# Number of worker threads for video frame processing
# Default: 4 workers (optimal for most systems)
# Set VIDEO_WORKERS env variable to override
MAX_VIDEO_WORKERS = int(os.getenv('VIDEO_WORKERS', '4'))

# Scene detection configuration
# Threshold for detecting scene changes (0-1, higher = more sensitive)
# Default: 0.3 (detects significant scene changes)
# Set SCENE_THRESHOLD env variable to override
SCENE_THRESHOLD = float(os.getenv('SCENE_THRESHOLD', '0.3'))

# Minimum frames between scene changes to avoid detecting flickering as scenes
# Default: 15 frames (0.5 seconds at 30fps)
# Set MIN_SCENE_DURATION env variable to override
MIN_SCENE_DURATION = int(os.getenv('MIN_SCENE_DURATION', '15'))

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

app = FastAPI(title="Avinash-EYE Multi-Media AI Service")

# Global variables for models
blip_processor = None
blip_model = None
clip_processor = None
clip_model = None
florence_processor = None
florence_model = None
siglip_processor = None
siglip_model = None
aimv2_processor = None
aimv2_model = None
whisper_model = None
device = None


def extract_json_from_response(text: str) -> Optional[dict]:
    """
    Extract JSON from LLM response text using balanced brace algorithm.
    Handles nested objects/arrays and multiple JSON blocks.

    Args:
        text: Raw text response from LLM that may contain JSON

    Returns:
        Parsed dict if valid JSON found, None otherwise
    """
    import re

    if not text or not isinstance(text, str):
        return None

    # Strategy 1: Try parsing the entire text as JSON first
    try:
        return json.loads(text.strip())
    except json.JSONDecodeError:
        pass

    # Strategy 2: Find JSON using balanced brace counting
    # This handles nested objects like {"a": {"b": 1}, "c": [1, 2]}
    start_idx = text.find('{')
    if start_idx == -1:
        return None

    brace_count = 0
    in_string = False
    escape_next = False

    for i, char in enumerate(text[start_idx:], start=start_idx):
        if escape_next:
            escape_next = False
            continue

        if char == '\\' and in_string:
            escape_next = True
            continue

        if char == '"' and not escape_next:
            in_string = not in_string
            continue

        if in_string:
            continue

        if char == '{':
            brace_count += 1
        elif char == '}':
            brace_count -= 1
            if brace_count == 0:
                # Found complete JSON object
                json_str = text[start_idx:i+1]
                try:
                    return json.loads(json_str)
                except json.JSONDecodeError:
                    # Try finding next JSON block
                    remaining = text[i+1:]
                    return extract_json_from_response(remaining)

    # Strategy 3: Try regex for simple cases (fallback)
    # More permissive pattern for edge cases
    simple_match = re.search(r'\{[^{}]+\}', text)
    if simple_match:
        try:
            return json.loads(simple_match.group())
        except json.JSONDecodeError:
            pass

    return None


# Request/Response Models
class AnalyzeImageRequest(BaseModel):
    """Request model for image analysis."""
    image_path: str
    use_ollama: bool = True
    detect_faces: bool = True
    captioning_model: str = "florence"  # "blip" or "florence" (florence recommended)
    embedding_model: str = "aimv2"  # "clip", "siglip", or "aimv2" (aimv2 recommended - outperforms SigLIP)
    ollama_model: str = "llava:13b-v1.6"  # Ollama model for detailed descriptions
    # Maximum analysis coverage options
    detect_objects: bool = True  # Object detection via Florence-2 <OD>
    extract_colors: bool = True  # Dominant color extraction via K-means
    analyze_quality: bool = True  # Image quality metrics via OpenCV
    compute_hashes: bool = True  # pHash/dHash for duplicate detection
    classify_scene: bool = True  # Scene classification via Ollama


class AnalyzeVideoRequest(BaseModel):
    """Request model for video analysis."""
    video_path: str
    extract_frames: bool = True
    frame_interval: int = 30  # Extract 1 frame every N frames


class AnalyzeDocumentRequest(BaseModel):
    """Request model for document analysis."""
    document_path: str
    perform_ocr: bool = True
    ocr_engine: str = "auto"  # "auto", "paddleocr", "tesseract"
    use_ollama: bool = False
    ollama_model: str = "qwen2.5:7b"


class TranscribeAudioRequest(BaseModel):
    """Request model for audio transcription."""
    audio_path: str
    language: Optional[str] = None  # Auto-detect if None


class EmbedTextRequest(BaseModel):
    """Request model for text embedding."""
    query: str


class ExtractEmailRequest(BaseModel):
    """Request model for email extraction."""
    file_path: str


class ExtractArchiveMetadataRequest(BaseModel):
    """Request model for archive metadata extraction."""
    file_path: str


class AnalyzeCodeFileRequest(BaseModel):
    """Request model for code file analysis."""
    file_path: str


class AnalyzeImageComprehensiveRequest(BaseModel):
    """Request model for comprehensive 4-pass VLM image analysis."""
    image_path: str
    analysis_mode: str = "comprehensive"  # "comprehensive" (4-pass ~40-60s) or "quick" (1-pass ~8s)
    ollama_model: str = "llava:13b-v1.6"  # Vision model for all passes
    embedding_model: str = "aimv2"  # Embedding model for similarity search
    detect_faces: bool = True
    generate_embedding: bool = True


class AnalyzeImageComprehensiveResponse(BaseModel):
    """Response model for comprehensive 4-pass VLM image analysis."""
    success: bool
    analysis_mode: str  # "comprehensive" or "quick"
    # 4-pass analysis results (JSONB-ready)
    content_analysis: Optional[Dict[str, Any]] = None  # Pass 1: Content & Scene
    people_analysis: Optional[Dict[str, Any]] = None  # Pass 2: People & Emotion
    quality_analysis: Optional[Dict[str, Any]] = None  # Pass 3: Quality & Technical
    context_analysis: Optional[Dict[str, Any]] = None  # Pass 4: Context & Metadata
    combined_analysis: Optional[Dict[str, Any]] = None  # Merged results with summary
    # Analysis metadata
    passes_completed: int = 0
    passes_failed: int = 0
    analysis_duration_seconds: float = 0.0
    errors: List[str] = []
    # Standard image analysis fields
    embedding: Optional[List[float]] = None
    faces_detected: int = 0
    face_locations: List[List[int]] = []
    face_encodings: List[List[float]] = []
    thumbnail_path: Optional[str] = None


class AnalyzeImageResponse(BaseModel):
    """Response model for image analysis."""
    description: str
    detailed_description: Optional[str] = None
    meta_tags: List[str] = []
    embedding: Optional[List[float]] = None  # None for non-raster images like SVG
    faces_detected: int = 0
    face_locations: List[List[int]] = []
    face_encodings: List[List[float]] = []
    thumbnail_path: Optional[str] = None  # Path to generated thumbnail (JPEG)
    extracted_text: Optional[str] = ""  # Extracted text for SVG/other special formats
    # Maximum analysis coverage fields
    objects_detected: Optional[Dict[str, Any]] = None  # Florence-2 <OD> results
    scene_classification: Optional[Dict[str, Any]] = None  # Ollama scene classification
    dominant_colors: Optional[List[Dict[str, Any]]] = None  # K-means color extraction
    image_quality: Optional[Dict[str, Any]] = None  # OpenCV quality metrics
    quality_tier: Optional[str] = None  # excellent, good, fair, poor
    phash: Optional[str] = None  # Perceptual hash
    dhash: Optional[str] = None  # Difference hash


class AnalyzeVideoResponse(BaseModel):
    """Response model for video analysis."""
    duration_seconds: float
    frame_count: int
    fps: float
    resolution: str
    scene_descriptions: List[Dict[str, Any]] = []
    embedding: List[float]  # Average embedding of key frames
    objects_detected: List[str] = []
    thumbnail_path: Optional[str] = None  # Path to generated thumbnail (JPEG)


class AnalyzeDocumentResponse(BaseModel):
    """Response model for document analysis."""
    extracted_text: str
    page_count: Optional[int] = None
    summary: Optional[str] = None
    keywords: List[str] = []
    embedding: List[float]
    thumbnail_path: Optional[str] = None  # Path to generated thumbnail (JPEG)
    # Intelligent document analysis fields
    document_type: Optional[str] = None  # Classified document type (invoice, receipt, etc.)
    classification_confidence: Optional[float] = None  # Confidence score (0.0-1.0)
    entities: Optional[Dict[str, Any]] = None  # Extracted entities (dates, amounts, parties, etc.)


class TranscribeAudioResponse(BaseModel):
    """Response model for audio transcription."""
    text: str
    language: str
    confidence: float
    embedding: List[float]
    thumbnail_path: Optional[str] = None  # Path to generated waveform thumbnail (JPEG)


class EmbedTextResponse(BaseModel):
    """Response model for text embedding."""
    embedding: List[float]


class ExtractEmailResponse(BaseModel):
    """Response model for email extraction."""
    sender: Optional[str] = None
    recipients: List[str] = []
    subject: Optional[str] = None
    date: Optional[str] = None
    body: str = ""
    attachment_count: int = 0
    has_html: bool = False


class ExtractArchiveMetadataResponse(BaseModel):
    """Response model for archive metadata extraction."""
    file_count: int = 0
    total_size: int = 0
    file_types: Dict[str, int] = {}
    file_list: List[Dict[str, Any]] = []


class AnalyzeCodeFileResponse(BaseModel):
    """Response model for code file analysis."""
    language: str = "unknown"
    line_count: int = 0
    code_lines: int = 0
    comment_lines: int = 0
    blank_lines: int = 0
    file_size: int = 0
    encoding: str = "utf-8"
    extracted_text: str = ""  # Full text content for searchability


@app.on_event("startup")
async def load_models():
    """Load AI models on startup."""
    global blip_processor, blip_model, clip_processor, clip_model, whisper_model, device

    logger.info("Starting multi-media model loading process...")

    # Determine device
    device = torch.device("cuda" if torch.cuda.is_available() else "cpu")
    logger.info(f"Using device: {device}")

    try:
        # Load BLIP model for image captioning
        logger.info("Loading BLIP model...")
        blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-large")
        blip_model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-large")
        blip_model.to(device)
        blip_model.eval()
        logger.info("BLIP model loaded successfully!")

        # Load CLIP model for embeddings
        logger.info("Loading CLIP model...")
        clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
        clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
        clip_model.to(device)
        clip_model.eval()
        logger.info("CLIP model loaded successfully!")

        # Load Whisper model for audio transcription
        if WHISPER_AVAILABLE:
            logger.info("Loading Whisper model...")
            whisper_model = whisper.load_model("base")
            logger.info("Whisper model loaded successfully!")

        logger.info("All models loaded and ready!")

    except Exception as e:
        logger.error(f"Error loading models: {str(e)}")
        raise


# ===== IMAGE PROCESSING =====

def generate_caption_blip(image: Image.Image) -> str:
    """Generate caption using BLIP."""
    inputs = blip_processor(image, return_tensors="pt").to(device)

    with torch.no_grad():
        out = blip_model.generate(
            **inputs,
            max_length=150,
            num_beams=5,
            temperature=1.0
        )

    return blip_processor.decode(out[0], skip_special_tokens=True)


def get_florence_model():
    """Lazily initialize Florence-2 model (it's heavy to load)."""
    global florence_processor, florence_model
    if florence_processor is None or florence_model is None:
        logger.info("Loading Florence-2 model...")
        try:
            # Use Florence-2-base for better performance on Apple Silicon
            model_name = "microsoft/Florence-2-base"
            florence_processor = AutoProcessor.from_pretrained(model_name, trust_remote_code=True)
            florence_model = AutoModelForCausalLM.from_pretrained(
                model_name,
                trust_remote_code=True,
                torch_dtype=torch.float32  # Use float32 for CPU/Apple Silicon compatibility
            )
            florence_model.to(device)
            florence_model.eval()
            logger.info("Florence-2 model loaded successfully!")
        except Exception as e:
            logger.error(f"Failed to load Florence-2: {str(e)}")
            return None, None
    return florence_processor, florence_model


def generate_caption_florence(image: Image.Image, detailed: bool = True) -> str:
    """Generate caption using Florence-2."""
    processor, model = get_florence_model()
    if processor is None or model is None:
        logger.warning("Florence-2 not available, falling back to BLIP")
        return generate_caption_blip(image)

    try:
        # Use DETAILED_CAPTION for more comprehensive descriptions
        task_prompt = "<MORE_DETAILED_CAPTION>" if detailed else "<CAPTION>"

        inputs = processor(text=task_prompt, images=image, return_tensors="pt").to(device)

        with torch.no_grad():
            generated_ids = model.generate(
                input_ids=inputs["input_ids"],
                pixel_values=inputs["pixel_values"],
                max_new_tokens=256,
                num_beams=3,
                do_sample=False
            )

        generated_text = processor.batch_decode(generated_ids, skip_special_tokens=False)[0]

        # Parse Florence-2 output format
        parsed = processor.post_process_generation(
            generated_text,
            task=task_prompt,
            image_size=(image.width, image.height)
        )

        # Extract the caption from the parsed result
        if isinstance(parsed, dict):
            caption = parsed.get(task_prompt, generated_text)
        else:
            caption = str(parsed)

        return caption.strip()

    except Exception as e:
        logger.error(f"Florence-2 caption generation failed: {str(e)}")
        return generate_caption_blip(image)


def generate_caption(image: Image.Image, model: str = "blip") -> str:
    """Generate caption using the specified model."""
    if model.lower() == "florence" or model.lower() == "florence-2":
        return generate_caption_florence(image, detailed=True)
    else:
        return generate_caption_blip(image)


def detect_faces(image: Image.Image) -> Dict:
    """Detect faces in image and return locations and encodings."""
    try:
        img_array = np.array(image)
        face_locations = face_recognition.face_locations(img_array)

        face_encodings = []
        if face_locations:
            face_encodings = face_recognition.face_encodings(img_array, face_locations)

        return {
            "count": len(face_locations),
            "locations": face_locations,
            "encodings": [encoding.tolist() for encoding in face_encodings]
        }
    except Exception as e:
        logger.error(f"Face detection failed: {str(e)}")
        return {"count": 0, "locations": [], "encodings": []}


# ============================================================================
# Maximum Analysis Coverage Functions
# ============================================================================

def detect_objects_florence(image: Image.Image) -> Dict:
    """
    Detect objects in image using Florence-2 <OD> task.
    Zero additional memory cost - reuses the already loaded Florence-2 model.

    Returns:
        Dict with labels, bboxes, and label_counts
    """
    processor, model = get_florence_model()
    if processor is None or model is None:
        logger.warning("Florence-2 not available for object detection")
        return {"labels": [], "bboxes": [], "label_counts": {}}

    try:
        task_prompt = "<OD>"  # Object Detection task

        inputs = processor(text=task_prompt, images=image, return_tensors="pt").to(device)

        with torch.no_grad():
            generated_ids = model.generate(
                input_ids=inputs["input_ids"],
                pixel_values=inputs["pixel_values"],
                max_new_tokens=1024,
                num_beams=3,
                do_sample=False
            )

        generated_text = processor.batch_decode(generated_ids, skip_special_tokens=False)[0]

        # Parse Florence-2 output format
        parsed = processor.post_process_generation(
            generated_text,
            task=task_prompt,
            image_size=(image.width, image.height)
        )

        # Extract labels and bboxes with robust multi-format parsing
        labels = []
        bboxes = []

        # Florence-2 can return different output formats depending on version/config
        # Strategy 1: Standard format with task key {'<OD>': {'bboxes': [], 'labels': []}}
        if isinstance(parsed, dict) and task_prompt in parsed:
            od_result = parsed[task_prompt]
            if isinstance(od_result, dict):
                labels = od_result.get("labels", [])
                bboxes = od_result.get("bboxes", [])
                logger.debug(f"Parsed using standard format (task key present)")

        # Strategy 2: Flat format without task key {'bboxes': [], 'labels': []}
        if not labels and isinstance(parsed, dict):
            if "labels" in parsed or "bboxes" in parsed:
                labels = parsed.get("labels", [])
                bboxes = parsed.get("bboxes", [])
                logger.debug(f"Parsed using flat format (no task key)")

        # Strategy 3: Check for 'objects' key (alternative format)
        if not labels and isinstance(parsed, dict) and "objects" in parsed:
            objects = parsed.get("objects", [])
            if isinstance(objects, list):
                for obj in objects:
                    if isinstance(obj, dict):
                        if "label" in obj:
                            labels.append(obj["label"])
                        if "bbox" in obj:
                            bboxes.append(obj["bbox"])
                logger.debug(f"Parsed using objects array format")

        # Strategy 4: Parse raw generated_text as fallback using robust JSON extraction
        if not labels and generated_text:
            fallback_parsed = extract_json_from_response(generated_text)
            if fallback_parsed and isinstance(fallback_parsed, dict):
                labels = fallback_parsed.get("labels", [])
                bboxes = fallback_parsed.get("bboxes", [])
                if labels:
                    logger.debug(f"Parsed using JSON extraction from generated_text")

        # Log warning if still empty after all strategies
        if not labels:
            logger.warning(f"Object detection returned no labels. Raw output type: {type(parsed)}, keys: {parsed.keys() if isinstance(parsed, dict) else 'N/A'}")

        # Count label occurrences
        from collections import Counter
        label_counts = dict(Counter(labels))

        logger.info(f"Object detection found {len(labels)} objects: {label_counts}")

        return {
            "labels": labels,
            "bboxes": [[int(b) for b in bbox] for bbox in bboxes] if bboxes else [],
            "label_counts": label_counts
        }

    except Exception as e:
        logger.error(f"Florence-2 object detection failed: {str(e)}")
        return {"labels": [], "bboxes": [], "label_counts": {}}


def extract_dominant_colors(image: Image.Image, n_colors: int = 5) -> List[Dict]:
    """
    Extract dominant colors from image using K-means clustering.
    CPU-only, memory efficient (~50MB for processing).

    Returns:
        List of dicts with hex, rgb, name, and percentage for each dominant color
    """
    try:
        from sklearn.cluster import KMeans

        # Resize image for faster processing (max 200x200)
        img = image.copy()
        img.thumbnail((200, 200), Image.Resampling.LANCZOS)

        # Convert to numpy array
        img_array = np.array(img)

        # Handle grayscale images
        if len(img_array.shape) == 2:
            img_array = np.stack([img_array] * 3, axis=-1)
        elif img_array.shape[2] == 4:  # RGBA
            img_array = img_array[:, :, :3]

        # Reshape to (n_pixels, 3)
        pixels = img_array.reshape(-1, 3)

        # Filter out very dark and very light pixels for better color detection
        brightness = np.mean(pixels, axis=1)
        mask = (brightness > 20) & (brightness < 235)
        filtered_pixels = pixels[mask] if np.sum(mask) > n_colors else pixels

        # K-means clustering
        kmeans = KMeans(n_clusters=min(n_colors, len(filtered_pixels)), random_state=42, n_init=10)
        kmeans.fit(filtered_pixels)

        # Get cluster centers and counts
        centers = kmeans.cluster_centers_.astype(int)
        labels = kmeans.labels_

        # Calculate percentage for each color
        from collections import Counter
        label_counts = Counter(labels)
        total = len(labels)

        colors = []
        for i, center in enumerate(centers):
            r, g, b = center
            hex_color = f"#{r:02x}{g:02x}{b:02x}"
            percentage = (label_counts[i] / total) * 100

            colors.append({
                "hex": hex_color,
                "rgb": [int(r), int(g), int(b)],
                "name": _get_color_name(r, g, b),
                "percentage": round(percentage, 1)
            })

        # Sort by percentage (most dominant first)
        colors.sort(key=lambda x: x["percentage"], reverse=True)

        logger.info(f"Extracted {len(colors)} dominant colors")
        return colors

    except Exception as e:
        logger.error(f"Color extraction failed: {str(e)}")
        return []


def _get_color_name(r: int, g: int, b: int) -> str:
    """Get approximate color name from RGB values."""
    # Simple color naming based on hue and saturation
    max_c = max(r, g, b)
    min_c = min(r, g, b)
    diff = max_c - min_c

    # Grayscale detection
    if diff < 20:
        if max_c < 50:
            return "black"
        elif max_c < 128:
            return "gray"
        elif max_c < 200:
            return "silver"
        else:
            return "white"

    # Color detection based on dominant channel
    if r >= g and r >= b:
        if g > b + 30:
            return "orange" if g > 128 else "brown"
        elif b > g + 30:
            return "pink" if r > 180 else "purple"
        else:
            return "red"
    elif g >= r and g >= b:
        if b > r + 30:
            return "cyan" if b > 128 else "teal"
        elif r > b + 30:
            return "yellow" if r > 180 else "olive"
        else:
            return "green"
    else:  # b is dominant
        if r > g + 30:
            return "purple" if r > 100 else "indigo"
        elif g > r + 30:
            return "cyan"
        else:
            return "blue"


def analyze_image_quality(image: Image.Image) -> Dict:
    """
    Analyze image quality using OpenCV metrics.

    Returns:
        Dict with overall_score, sharpness, brightness, contrast, saturation, noise, and issues
    """
    try:
        # Convert PIL to OpenCV format
        img_array = np.array(image)
        if len(img_array.shape) == 2:
            gray = img_array
            img_bgr = cv2.cvtColor(img_array, cv2.COLOR_GRAY2BGR)
        else:
            img_bgr = cv2.cvtColor(img_array, cv2.COLOR_RGB2BGR)
            gray = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2GRAY)

        # 1. Sharpness (Laplacian variance)
        laplacian_var = cv2.Laplacian(gray, cv2.CV_64F).var()
        # Normalize: <100 is blurry, >500 is very sharp
        sharpness = min(1.0, laplacian_var / 500)
        is_blurry = laplacian_var < 100

        # 2. Brightness (mean luminance)
        brightness_raw = np.mean(gray)
        # Ideal brightness is around 128 (middle gray)
        brightness = 1.0 - abs(brightness_raw - 128) / 128
        is_dark = brightness_raw < 50
        is_overexposed = brightness_raw > 220

        # 3. Contrast (standard deviation of luminance)
        contrast_raw = np.std(gray)
        # Normalize: <30 is low contrast, >70 is good contrast
        contrast = min(1.0, contrast_raw / 70)
        is_low_contrast = contrast_raw < 30

        # 4. Saturation (for color images)
        saturation = 0.5  # default for grayscale
        is_desaturated = False
        if len(img_array.shape) == 3 and img_array.shape[2] >= 3:
            hsv = cv2.cvtColor(img_bgr, cv2.COLOR_BGR2HSV)
            saturation_raw = np.mean(hsv[:, :, 1])
            saturation = saturation_raw / 255
            is_desaturated = saturation_raw < 30

        # 5. Noise estimation (high-frequency content in smooth areas)
        # Use median filter comparison
        denoised = cv2.medianBlur(gray, 5)
        noise_raw = np.mean(np.abs(gray.astype(float) - denoised.astype(float)))
        # Normalize: <5 is clean, >20 is noisy
        noise_score = max(0, 1.0 - noise_raw / 20)
        is_noisy = noise_raw > 15

        # Calculate overall score (weighted average)
        overall_score = (
            sharpness * 0.35 +
            brightness * 0.20 +
            contrast * 0.20 +
            saturation * 0.10 +
            noise_score * 0.15
        )

        # Determine quality tier
        if overall_score >= 0.8:
            quality_tier = "excellent"
        elif overall_score >= 0.6:
            quality_tier = "good"
        elif overall_score >= 0.4:
            quality_tier = "fair"
        else:
            quality_tier = "poor"

        # Convert numpy bools to Python native bools for JSON serialization
        issues = {
            "is_blurry": bool(is_blurry),
            "is_dark": bool(is_dark),
            "is_overexposed": bool(is_overexposed),
            "is_low_contrast": bool(is_low_contrast),
            "is_desaturated": bool(is_desaturated),
            "is_noisy": bool(is_noisy)
        }

        result = {
            "overall_score": round(overall_score, 2),
            "sharpness": round(sharpness, 2),
            "brightness": round(brightness, 2),
            "contrast": round(contrast, 2),
            "saturation": round(saturation, 2),
            "noise_score": round(noise_score, 2),
            "quality_tier": quality_tier,
            "issues": issues
        }

        logger.info(f"Image quality: {quality_tier} (score: {overall_score:.2f})")
        return result

    except Exception as e:
        logger.error(f"Image quality analysis failed: {str(e)}")
        return {
            "overall_score": 0.5,
            "sharpness": 0.5,
            "brightness": 0.5,
            "contrast": 0.5,
            "saturation": 0.5,
            "noise_score": 0.5,
            "quality_tier": "unknown",
            "issues": {}
        }


def compute_perceptual_hashes(image: Image.Image) -> Dict[str, str]:
    """
    Compute perceptual hashes for duplicate detection.

    pHash (perceptual hash) - best for similar images with different sizes/formats
    dHash (difference hash) - good for detecting crops and edits

    Returns:
        Dict with phash and dhash as hex strings
    """
    try:
        import imagehash

        # Ensure image is in a format imagehash can process
        img = image.copy()
        if img.mode not in ('RGB', 'L'):
            img = img.convert('RGB')

        # pHash - perceptual hash using DCT
        phash = str(imagehash.phash(img))

        # dHash - difference hash (horizontal gradient)
        dhash = str(imagehash.dhash(img))

        logger.info(f"Computed hashes: pHash={phash}, dHash={dhash}")

        return {
            "phash": phash,
            "dhash": dhash
        }

    except ImportError:
        logger.error("imagehash library not available")
        return {"phash": None, "dhash": None}
    except Exception as e:
        logger.error(f"Hash computation failed: {str(e)}")
        return {"phash": None, "dhash": None}


def classify_scene_ollama(image: Image.Image, ollama_model: str = "llava:13b-v1.6") -> Dict:
    """
    Classify scene/environment using Ollama vision model.

    Tries requested model first, falls back to llava:latest if not available.

    Returns:
        Dict with environment, setting, context, mood, and confidence
    """
    if not OLLAMA_AVAILABLE or ollama_client is None:
        logger.warning("Ollama not available for scene classification")
        return {}

    # Models to try in order of preference
    models_to_try = [ollama_model]
    if ollama_model != "llava:latest":
        models_to_try.append("llava:latest")

    try:
        import base64
        from io import BytesIO

        # Convert image to base64
        img = image.copy()
        img.thumbnail((1024, 1024), Image.Resampling.LANCZOS)  # Limit size for Ollama

        buffer = BytesIO()
        img.save(buffer, format="JPEG", quality=85)
        img_base64 = base64.b64encode(buffer.getvalue()).decode('utf-8')

        prompt = """Analyze this image and provide a structured scene classification. Return ONLY a JSON object with these fields:
{
    "environment": "indoor" or "outdoor" or "mixed",
    "setting": brief location type (e.g., "beach", "office", "forest", "city street", "living room"),
    "context": list of 2-4 contextual tags (e.g., ["vacation", "relaxation"], ["work", "professional"]),
    "mood": overall mood/atmosphere (e.g., "peaceful", "energetic", "formal", "casual"),
    "time_of_day": if determinable ("morning", "afternoon", "evening", "night", "unknown"),
    "weather": if outdoor and determinable ("sunny", "cloudy", "rainy", "unknown"),
    "confidence": your confidence in this classification (0.0 to 1.0)
}
Return ONLY valid JSON, no other text."""

        last_error = None
        for model in models_to_try:
            try:
                logger.info(f"Attempting scene classification with model: {model}")
                response = ollama_client.generate(
                    model=model,
                    prompt=prompt,
                    images=[img_base64],
                    options={"temperature": 0.3}  # Lower temperature for more consistent output
                )

                # Parse JSON response using robust extraction
                response_text = response.get('response', '').strip()
                result = extract_json_from_response(response_text)

                if result:
                    logger.info(f"Scene classified with {model}: {result.get('setting', 'unknown')} ({result.get('environment', 'unknown')})")
                    result['model_used'] = model  # Track which model was used
                    return result
                else:
                    logger.warning(f"Could not parse scene classification response from {model}: {response_text[:100]}")
                    last_error = f"Failed to parse response from {model}"
                    continue

            except Exception as e:
                error_msg = str(e)
                if "not found" in error_msg.lower() or "404" in error_msg:
                    logger.warning(f"Model {model} not available, trying next...")
                    last_error = f"Model {model} not found"
                    continue
                else:
                    logger.error(f"Scene classification with {model} failed: {error_msg}")
                    last_error = error_msg
                    continue

        # All models failed
        logger.error(f"Scene classification failed with all models. Last error: {last_error}")
        return {}

    except Exception as e:
        logger.error(f"Scene classification failed: {str(e)}")
        return {}


def generate_image_embedding_clip(image: Image.Image) -> np.ndarray:
    """Generate normalized embedding vector using CLIP."""
    inputs = clip_processor(images=image, return_tensors="pt").to(device)

    with torch.no_grad():
        image_features = clip_model.get_image_features(**inputs)

    embedding = image_features / image_features.norm(dim=-1, keepdim=True)
    return embedding.cpu().numpy().flatten()


def get_siglip_model():
    """Lazily initialize SigLIP model (it's heavy to load)."""
    global siglip_processor, siglip_model
    if siglip_processor is None or siglip_model is None:
        logger.info("Loading SigLIP model...")
        try:
            from transformers import AutoProcessor, AutoModel
            # Use SigLIP base for good balance of quality and speed
            model_name = "google/siglip-base-patch16-224"
            siglip_processor = AutoProcessor.from_pretrained(model_name)
            siglip_model = AutoModel.from_pretrained(model_name)
            siglip_model.to(device)
            siglip_model.eval()
            logger.info("SigLIP model loaded successfully!")
        except Exception as e:
            logger.error(f"Failed to load SigLIP: {str(e)}")
            return None, None
    return siglip_processor, siglip_model


def generate_image_embedding_siglip(image: Image.Image) -> np.ndarray:
    """Generate normalized embedding vector using SigLIP."""
    processor, model = get_siglip_model()
    if processor is None or model is None:
        logger.warning("SigLIP not available, falling back to CLIP")
        return generate_image_embedding_clip(image)

    try:
        inputs = processor(images=image, return_tensors="pt").to(device)

        with torch.no_grad():
            outputs = model.get_image_features(**inputs)

        embedding = outputs / outputs.norm(dim=-1, keepdim=True)
        return embedding.cpu().numpy().flatten()

    except Exception as e:
        logger.error(f"SigLIP embedding failed: {str(e)}")
        return generate_image_embedding_clip(image)


def get_aimv2_model():
    """
    Lazily initialize AIMv2 model (Apple's Autoregressive Image Models v2).

    AIMv2 consistently outperforms state-of-the-art contrastive models like CLIP
    and SigLIP in multimodal image understanding across diverse settings.

    Available variants:
    - apple/aimv2-large-patch14-224 (recommended - good balance)
    - apple/aimv2-huge-patch14-224 (highest quality, more memory)
    - apple/aimv2-large-patch14-336 (larger input, more detail)
    """
    global aimv2_processor, aimv2_model
    if aimv2_processor is None or aimv2_model is None:
        logger.info("Loading AIMv2 model (Apple's state-of-the-art embedding model)...")
        try:
            from transformers import AutoProcessor, AutoModel
            # Use AIMv2-large for good balance of quality and speed on Apple Silicon
            model_name = "apple/aimv2-large-patch14-224"
            aimv2_processor = AutoProcessor.from_pretrained(model_name, trust_remote_code=True)
            aimv2_model = AutoModel.from_pretrained(model_name, trust_remote_code=True)
            aimv2_model.to(device)
            aimv2_model.eval()
            logger.info("AIMv2 model loaded successfully!")
        except Exception as e:
            logger.error(f"Failed to load AIMv2: {str(e)}")
            return None, None
    return aimv2_processor, aimv2_model


def generate_image_embedding_aimv2(image: Image.Image) -> np.ndarray:
    """
    Generate normalized embedding vector using AIMv2 (Apple's model).

    AIMv2 outperforms CLIP and SigLIP for image understanding and retrieval.
    """
    processor, model = get_aimv2_model()
    if processor is None or model is None:
        logger.warning("AIMv2 not available, falling back to SigLIP")
        return generate_image_embedding_siglip(image)

    try:
        inputs = processor(images=image, return_tensors="pt").to(device)

        with torch.no_grad():
            outputs = model(inputs["pixel_values"])

        # AIMv2 returns features that need to be extracted
        # Use the pooled output or mean of last hidden state
        if hasattr(outputs, 'pooler_output') and outputs.pooler_output is not None:
            features = outputs.pooler_output
        else:
            # Use mean pooling of last hidden state
            features = outputs.last_hidden_state.mean(dim=1)

        # Normalize the embedding
        embedding = features / features.norm(dim=-1, keepdim=True)
        return embedding.cpu().numpy().flatten()

    except Exception as e:
        logger.error(f"AIMv2 embedding failed: {str(e)}, falling back to SigLIP")
        return generate_image_embedding_siglip(image)


def generate_image_embedding(image: Image.Image, model: str = "aimv2") -> np.ndarray:
    """
    Generate embedding using the specified model.

    Args:
        image: PIL Image to generate embedding for
        model: Embedding model to use:
            - "aimv2" (recommended - Apple's model, outperforms CLIP and SigLIP)
            - "siglip" (Google's model, good accuracy)
            - "clip" (OpenAI's original, fast but less accurate)

    Returns:
        Normalized embedding vector as numpy array
    """
    model_lower = model.lower()
    if model_lower == "aimv2":
        return generate_image_embedding_aimv2(image)
    elif model_lower == "siglip":
        return generate_image_embedding_siglip(image)
    else:
        return generate_image_embedding_clip(image)


def generate_thumbnail(image_path: str, max_size: tuple = (800, 800)) -> Optional[str]:
    """
    Generate a browser-compatible JPEG thumbnail for any image format.

    Args:
        image_path: Path to the original image file
        max_size: Maximum dimensions (width, height) for the thumbnail

    Returns:
        Path to the generated thumbnail or None if generation fails
    """
    try:
        # Parse the original path
        original_path = Path(image_path)

        # Create thumbnail directory structure
        # Original: /app/shared/images/abc.heic
        # Thumbnail: /app/shared/images/thumbnails/abc.jpg
        thumbnail_dir = original_path.parent / 'thumbnails'
        thumbnail_dir.mkdir(parents=True, exist_ok=True)

        # Generate thumbnail filename (always .jpg)
        thumbnail_filename = original_path.stem + '.jpg'
        thumbnail_path = thumbnail_dir / thumbnail_filename

        # Open and convert the image
        image = Image.open(image_path).convert("RGB")

        # Generate thumbnail (maintains aspect ratio)
        image.thumbnail(max_size, Image.Resampling.LANCZOS)

        # Save as JPEG
        image.save(str(thumbnail_path), "JPEG", quality=85, optimize=True)

        logger.info(f"Generated thumbnail: {thumbnail_path}")
        return str(thumbnail_path)

    except Exception as e:
        logger.error(f"Failed to generate thumbnail for {image_path}: {str(e)}")
        return None


# ===== VIDEO PROCESSING =====

def generate_video_thumbnail(video_path: str, time_position: float = 1.0, max_size: tuple = (800, 800)) -> Optional[str]:
    """
    Generate a thumbnail from a video file by extracting a frame.

    Args:
        video_path: Path to the video file
        time_position: Time position in seconds to extract frame (default: 1.0)
        max_size: Maximum dimensions for the thumbnail

    Returns:
        Path to the generated thumbnail or None if generation fails
    """
    try:
        # Parse the original path
        original_path = Path(video_path)

        # Create thumbnail directory
        thumbnail_dir = original_path.parent / 'thumbnails'
        thumbnail_dir.mkdir(parents=True, exist_ok=True)

        # Generate thumbnail filename
        thumbnail_filename = original_path.stem + '.jpg'
        thumbnail_path = thumbnail_dir / thumbnail_filename

        # Open video and get metadata
        cap = cv2.VideoCapture(video_path)
        if not cap.isOpened():
            logger.error(f"Could not open video file: {video_path}")
            return None

        fps = cap.get(cv2.CAP_PROP_FPS)
        total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        duration = total_frames / fps if fps > 0 else 0

        logger.info(f"Video info: duration={duration:.2f}s, fps={fps}, total_frames={total_frames}")

        # Try multiple time positions as fallbacks
        positions = [time_position, 0.5, 2.0, 5.0, 0.0]

        for pos in positions:
            # Skip positions beyond video duration
            if pos > duration:
                continue

            frame_number = min(int(fps * pos), total_frames - 1) if fps > 0 else 30
            cap.set(cv2.CAP_PROP_POS_FRAMES, frame_number)
            ret, frame = cap.read()

            if ret and frame is not None and frame.size > 0:
                # Success! Convert and save the frame
                try:
                    # Convert BGR to RGB
                    frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)

                    # Convert to PIL Image
                    image = Image.fromarray(frame_rgb)

                    # Generate thumbnail
                    image.thumbnail(max_size, Image.Resampling.LANCZOS)

                    # Save as JPEG
                    image.save(str(thumbnail_path), "JPEG", quality=85, optimize=True)

                    cap.release()
                    logger.info(f"Generated video thumbnail at {pos}s: {thumbnail_path}")
                    return str(thumbnail_path)
                except Exception as e:
                    logger.warning(f"Failed to process frame at {pos}s: {str(e)}")
                    continue

        cap.release()
        logger.error(f"Failed to extract valid frame from video after trying all positions: {video_path}")
        return None

    except Exception as e:
        logger.error(f"Failed to generate video thumbnail for {video_path}: {str(e)}")
        return None


def _calculate_frame_difference(frame1: np.ndarray, frame2: np.ndarray) -> float:
    """
    Calculate the difference between two frames using histogram comparison.

    Returns a value between 0 (identical) and 1 (completely different).
    Uses HSV color space for better scene detection accuracy.
    """
    try:
        # Convert to HSV for better color comparison
        hsv1 = cv2.cvtColor(frame1, cv2.COLOR_RGB2HSV)
        hsv2 = cv2.cvtColor(frame2, cv2.COLOR_RGB2HSV)

        # Calculate histograms
        hist1 = cv2.calcHist([hsv1], [0, 1], None, [50, 60], [0, 180, 0, 256])
        hist2 = cv2.calcHist([hsv2], [0, 1], None, [50, 60], [0, 180, 0, 256])

        # Normalize histograms
        cv2.normalize(hist1, hist1, alpha=0, beta=1, norm_type=cv2.NORM_MINMAX)
        cv2.normalize(hist2, hist2, alpha=0, beta=1, norm_type=cv2.NORM_MINMAX)

        # Compare histograms (returns value between 0 and 1)
        # 1 = identical, 0 = completely different
        similarity = cv2.compareHist(hist1, hist2, cv2.HISTCMP_CORREL)

        # Convert to difference (0 = identical, 1 = completely different)
        difference = 1.0 - similarity

        return difference

    except Exception as e:
        logger.error(f"Failed to calculate frame difference: {str(e)}")
        return 0.0


def extract_video_frames_with_scene_detection(
    video_path: str,
    scene_threshold: float = 0.3,
    min_scene_duration_frames: int = 15
) -> List[np.ndarray]:
    """
    Extract keyframes from video using smart scene detection.

    Only extracts frames when a significant scene change is detected, reducing
    redundant processing while capturing all important visual content.

    Args:
        video_path: Path to the video file
        scene_threshold: Threshold for scene change detection (0-1, higher = more sensitive)
        min_scene_duration_frames: Minimum frames between scene changes to avoid flickering

    Returns:
        List of numpy arrays containing only keyframes from unique scenes
    """
    try:
        cap = cv2.VideoCapture(video_path)
        frames = []
        prev_frame = None
        frame_count = 0
        last_scene_frame = 0

        logger.info(f"Starting smart scene detection with threshold={scene_threshold}")

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            # Convert BGR to RGB
            frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)

            # Always add the first frame
            if prev_frame is None:
                frames.append(frame_rgb)
                prev_frame = frame
                last_scene_frame = frame_count
                logger.info(f"Frame {frame_count}: Added first frame")
                frame_count += 1
                continue

            # Check if enough frames have passed since last scene
            frames_since_last_scene = frame_count - last_scene_frame

            if frames_since_last_scene >= min_scene_duration_frames:
                # Calculate difference from previous frame
                difference = _calculate_frame_difference(prev_frame, frame)

                # If difference exceeds threshold, it's a new scene
                if difference >= scene_threshold:
                    frames.append(frame_rgb)
                    last_scene_frame = frame_count
                    logger.info(f"Frame {frame_count}: Scene change detected (diff={difference:.3f})")

            prev_frame = frame
            frame_count += 1

        cap.release()
        logger.info(f"Smart scene detection complete: extracted {len(frames)} keyframes from {frame_count} total frames")
        return frames

    except Exception as e:
        logger.error(f"Failed to extract video frames with scene detection: {str(e)}")
        return []


def extract_video_frames(video_path: str, frame_interval: int = 30) -> List[np.ndarray]:
    """
    Extract frames from video at specified interval (legacy method).

    This method is kept for backward compatibility. For better performance,
    use extract_video_frames_with_scene_detection() instead.
    """
    try:
        cap = cv2.VideoCapture(video_path)
        frames = []
        frame_count = 0

        while True:
            ret, frame = cap.read()
            if not ret:
                break

            if frame_count % frame_interval == 0:
                # Convert BGR to RGB
                frame_rgb = cv2.cvtColor(frame, cv2.COLOR_BGR2RGB)
                frames.append(frame_rgb)

            frame_count += 1

        cap.release()
        return frames

    except Exception as e:
        logger.error(f"Failed to extract video frames: {str(e)}")
        return []


def get_video_metadata(video_path: str) -> Dict:
    """Get video metadata using OpenCV."""
    try:
        cap = cv2.VideoCapture(video_path)

        fps = cap.get(cv2.CAP_PROP_FPS)
        frame_count = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
        width = int(cap.get(cv2.CAP_PROP_FRAME_WIDTH))
        height = int(cap.get(cv2.CAP_PROP_FRAME_HEIGHT))
        duration = frame_count / fps if fps > 0 else 0

        cap.release()

        return {
            "duration_seconds": duration,
            "frame_count": frame_count,
            "fps": fps,
            "resolution": f"{width}x{height}"
        }

    except Exception as e:
        logger.error(f"Failed to get video metadata: {str(e)}")
        return {}


def _process_single_frame(idx: int, frame: np.ndarray) -> Optional[Dict]:
    """Process a single video frame. Helper function for parallel processing."""
    try:
        # Convert numpy array to PIL Image
        pil_image = Image.fromarray(frame)

        # Generate caption for this frame
        caption = generate_caption_blip(pil_image)

        return {
            "frame_index": idx,
            "description": caption
        }

    except Exception as e:
        logger.error(f"Failed to analyze frame {idx}: {str(e)}")
        return None


def analyze_video_scenes(frames: List[np.ndarray]) -> List[Dict]:
    """
    Analyze video frames and generate scene descriptions using parallel processing.

    Uses ThreadPoolExecutor to process multiple frames concurrently for better performance.
    Number of workers is configurable via MAX_VIDEO_WORKERS constant.
    """
    scene_descriptions = []

    # Use parallel processing if we have multiple frames
    if len(frames) > 1:
        logger.info(f"Processing {len(frames)} video frames with {MAX_VIDEO_WORKERS} workers")

        with ThreadPoolExecutor(max_workers=MAX_VIDEO_WORKERS) as executor:
            # Submit all frames for processing
            future_to_idx = {
                executor.submit(_process_single_frame, idx, frame): idx
                for idx, frame in enumerate(frames)
            }

            # Collect results as they complete
            for future in as_completed(future_to_idx):
                result = future.result()
                if result is not None:
                    scene_descriptions.append(result)

        # Sort by frame index to maintain order
        scene_descriptions.sort(key=lambda x: x['frame_index'])

    else:
        # Single frame - no need for parallel processing
        result = _process_single_frame(0, frames[0])
        if result is not None:
            scene_descriptions.append(result)

    logger.info(f"Successfully analyzed {len(scene_descriptions)}/{len(frames)} frames")
    return scene_descriptions


# ===== DOCUMENT PROCESSING =====

def generate_document_thumbnail(document_path: str, max_size: tuple = (800, 800)) -> Optional[str]:
    """
    Generate a thumbnail from a document file (PDF, text, etc.).

    Args:
        document_path: Path to the document file
        max_size: Maximum dimensions for the thumbnail

    Returns:
        Path to the generated thumbnail or None if generation fails
    """
    try:
        # Parse the original path
        original_path = Path(document_path)

        # Create thumbnail directory
        thumbnail_dir = original_path.parent / 'thumbnails'
        thumbnail_dir.mkdir(parents=True, exist_ok=True)

        # Generate thumbnail filename
        thumbnail_filename = original_path.stem + '.jpg'
        thumbnail_path = thumbnail_dir / thumbnail_filename

        # Check document type
        mime_type = original_path.suffix.lower()

        if mime_type == '.pdf':
            # Try to use pdf2image if available
            try:
                from pdf2image import convert_from_path
                images = convert_from_path(str(document_path), first_page=1, last_page=1)
                if images:
                    image = images[0]
                    # Generate thumbnail
                    image.thumbnail(max_size, Image.Resampling.LANCZOS)
                    image.save(str(thumbnail_path), "JPEG", quality=85, optimize=True)
                    logger.info(f"Generated PDF thumbnail: {thumbnail_path}")
                    return str(thumbnail_path)
            except ImportError:
                logger.warning("pdf2image not available, trying PyMuPDF")
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(document_path)
                    if len(doc) > 0:
                        page = doc[0]
                        # Render page to pixmap
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom for quality
                        # Convert to PIL Image
                        img_data = pix.tobytes("ppm")
                        image = Image.frombytes("RGB", [pix.width, pix.height], img_data)
                        # Generate thumbnail
                        image.thumbnail(max_size, Image.Resampling.LANCZOS)
                        image.save(str(thumbnail_path), "JPEG", quality=85, optimize=True)
                        doc.close()
                        logger.info(f"Generated PDF thumbnail with PyMuPDF: {thumbnail_path}")
                        return str(thumbnail_path)
                except ImportError:
                    logger.warning("PyMuPDF not available, cannot generate PDF thumbnail")
                    return None

        elif mime_type in ['.txt', '.csv', '.log', '.md']:
            # For text files, create a text preview image with actual content
            with open(document_path, 'r', encoding='utf-8', errors='ignore') as f:
                text_content = f.read(1000)  # First 1000 chars

            from PIL import ImageDraw, ImageFont

            # Create image with paper-like background
            img = Image.new('RGB', (800, 1000), color='#f5f5f0')
            draw = ImageDraw.Draw(img)

            # Try to use a monospace font, fallback to default
            try:
                # Try common monospace fonts
                font_size = 12
                font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSansMono.ttf", font_size)
            except:
                try:
                    font = ImageFont.truetype("/System/Library/Fonts/Courier.dfont", font_size)
                except:
                    font = ImageFont.load_default()

            # Add padding and format text
            padding = 20
            max_width = 800 - (2 * padding)
            y_position = padding

            # Split text into lines and wrap
            lines = []
            for line in text_content.split('\n'):
                if not line:
                    lines.append('')
                    continue
                # Simple word wrapping
                words = line.split(' ')
                current_line = ''
                for word in words:
                    test_line = current_line + (' ' if current_line else '') + word
                    # Approximate character width (more accurate would use textbbox)
                    if len(test_line) * 7 < max_width:
                        current_line = test_line
                    else:
                        if current_line:
                            lines.append(current_line)
                        current_line = word
                if current_line:
                    lines.append(current_line)

            # Draw lines (max 60 lines to fit in image)
            for line in lines[:60]:
                draw.text((padding, y_position), line, fill='#333333', font=font)
                y_position += 16

            # Add truncation indicator if needed
            if len(lines) > 60 or len(text_content) >= 1000:
                draw.text((padding, y_position + 10), "...", fill='#666666', font=font)

            # Generate thumbnail
            img.thumbnail(max_size, Image.Resampling.LANCZOS)
            img.save(str(thumbnail_path), "JPEG", quality=85)
            logger.info(f"Generated text document thumbnail: {thumbnail_path}")
            return str(thumbnail_path)

        else:
            logger.warning(f"Unsupported document type for thumbnail: {mime_type}")
            return None

    except Exception as e:
        logger.error(f"Failed to generate document thumbnail for {document_path}: {str(e)}")
        return None


# ===== OFFICE DOCUMENT EXTRACTION =====

def extract_word_document(document_path: str) -> str:
    """Extract text from Word documents (.docx, .doc, .odt, .rtf)."""
    try:
        from pathlib import Path
        doc_path = Path(document_path)
        extension = doc_path.suffix.lower()

        if extension in ['.docx', '.doc']:
            # Handle Word documents
            try:
                from docx import Document
                doc = Document(document_path)
                full_text = []
                for para in doc.paragraphs:
                    if para.text.strip():
                        full_text.append(para.text)
                return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract Word document: {str(e)}")
                return ""

        elif extension == '.odt':
            # Handle OpenDocument Text
            try:
                from odf import text, teletype
                from odf.opendocument import load
                textdoc = load(document_path)
                all_paragraphs = textdoc.getElementsByType(text.P)
                full_text = [teletype.extractText(p) for p in all_paragraphs if teletype.extractText(p).strip()]
                return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract ODT document: {str(e)}")
                return ""

        elif extension == '.rtf':
            # RTF files - try to read as plain text (basic extraction)
            try:
                with open(document_path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
                # Remove RTF control codes (basic cleanup)
                import re
                text = re.sub(r'\\[a-z]+\d*\s?', '', content)
                text = re.sub(r'[{}]', '', text)
                return text.strip()
            except Exception as e:
                logger.error(f"Failed to extract RTF document: {str(e)}")
                return ""

        return ""

    except Exception as e:
        logger.error(f"Error extracting Word document: {str(e)}")
        return ""


def extract_spreadsheet(document_path: str) -> str:
    """Extract text and data from spreadsheets (.xlsx, .xls, .ods, .csv)."""
    try:
        from pathlib import Path
        doc_path = Path(document_path)
        extension = doc_path.suffix.lower()

        if extension in ['.xlsx', '.xls']:
            # Handle Excel files
            try:
                from openpyxl import load_workbook
                wb = load_workbook(document_path, read_only=True, data_only=True)
                full_text = []

                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    full_text.append(f"=== Sheet: {sheet_name} ===")

                    # Extract headers (first row)
                    headers = []
                    for cell in sheet[1]:
                        if cell.value:
                            headers.append(str(cell.value))

                    if headers:
                        full_text.append("Headers: " + ", ".join(headers))

                    # Extract data (sample first 20 rows to avoid huge text dumps)
                    for idx, row in enumerate(sheet.iter_rows(min_row=2, max_row=21, values_only=True), 2):
                        row_values = [str(cell) for cell in row if cell is not None]
                        if row_values:
                            full_text.append(f"Row {idx}: " + " | ".join(row_values))

                wb.close()
                return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract Excel document: {str(e)}")
                return ""

        elif extension == '.ods':
            # Handle OpenDocument Spreadsheet
            try:
                from odf.opendocument import load
                from odf.table import Table, TableRow, TableCell
                from odf import teletype

                spreadsheet = load(document_path)
                full_text = []

                tables = spreadsheet.getElementsByType(Table)
                for table in tables[:3]:  # Limit to first 3 sheets
                    table_name = table.getAttribute("name") or "Sheet"
                    full_text.append(f"=== Sheet: {table_name} ===")

                    rows = table.getElementsByType(TableRow)
                    for idx, row in enumerate(rows[:21], 1):  # First 20 rows
                        cells = row.getElementsByType(TableCell)
                        row_values = [teletype.extractText(cell).strip() for cell in cells if teletype.extractText(cell).strip()]
                        if row_values:
                            full_text.append(f"Row {idx}: " + " | ".join(row_values))

                return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract ODS document: {str(e)}")
                return ""

        elif extension == '.csv':
            # Handle CSV files
            try:
                import csv
                with open(document_path, 'r', encoding='utf-8', errors='ignore') as f:
                    reader = csv.reader(f)
                    rows = list(reader)[:21]  # First 20 rows
                    full_text = [" | ".join(row) for row in rows]
                    return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract CSV document: {str(e)}")
                return ""

        return ""

    except Exception as e:
        logger.error(f"Error extracting spreadsheet: {str(e)}")
        return ""


def extract_presentation(document_path: str) -> str:
    """Extract text from presentations (.pptx, .ppt, .odp)."""
    try:
        from pathlib import Path
        doc_path = Path(document_path)
        extension = doc_path.suffix.lower()

        if extension in ['.pptx', '.ppt']:
            # Handle PowerPoint files
            try:
                from pptx import Presentation
                prs = Presentation(document_path)
                full_text = []

                for slide_num, slide in enumerate(prs.slides, 1):
                    full_text.append(f"=== Slide {slide_num} ===")

                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            full_text.append(shape.text)

                    # Extract notes
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
                        full_text.append(f"Notes: {slide.notes_slide.notes_text_frame.text}")

                return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract PowerPoint document: {str(e)}")
                return ""

        elif extension == '.odp':
            # Handle OpenDocument Presentation
            try:
                from odf.opendocument import load
                from odf.text import P
                from odf import teletype

                presentation = load(document_path)
                full_text = []

                # Extract all text paragraphs
                paragraphs = presentation.getElementsByType(P)
                for para in paragraphs:
                    text = teletype.extractText(para).strip()
                    if text:
                        full_text.append(text)

                return '\n'.join(full_text)
            except Exception as e:
                logger.error(f"Failed to extract ODP document: {str(e)}")
                return ""

        return ""

    except Exception as e:
        logger.error(f"Error extracting presentation: {str(e)}")
        return ""


def extract_email(email_path: str) -> dict:
    """
    Extract metadata and content from email files (.eml, .msg).

    Returns:
        dict with sender, recipients, subject, date, body, and attachment count
    """
    try:
        from pathlib import Path
        import email
        from email import policy
        from email.parser import BytesParser
        import datetime

        email_file_path = Path(email_path)
        extension = email_file_path.suffix.lower()

        result = {
            'sender': None,
            'recipients': [],
            'subject': None,
            'date': None,
            'body': '',
            'attachment_count': 0,
            'has_html': False
        }

        if extension == '.eml':
            # Handle .eml files (standard RFC 822 format)
            try:
                with open(email_path, 'rb') as f:
                    msg = BytesParser(policy=policy.default).parse(f)

                # Extract sender
                result['sender'] = str(msg.get('From', ''))

                # Extract recipients
                to_addrs = msg.get('To', '')
                cc_addrs = msg.get('Cc', '')
                recipients = []
                if to_addrs:
                    recipients.extend([addr.strip() for addr in str(to_addrs).split(',')])
                if cc_addrs:
                    recipients.extend([addr.strip() for addr in str(cc_addrs).split(',')])
                result['recipients'] = recipients

                # Extract subject
                result['subject'] = str(msg.get('Subject', ''))

                # Extract date
                date_str = msg.get('Date')
                if date_str:
                    result['date'] = str(date_str)

                # Extract body
                body_parts = []
                if msg.is_multipart():
                    for part in msg.walk():
                        content_type = part.get_content_type()
                        if content_type == 'text/plain':
                            try:
                                body_parts.append(part.get_content())
                            except:
                                pass
                        elif content_type == 'text/html':
                            result['has_html'] = True

                        # Count attachments
                        if part.get_content_disposition() == 'attachment':
                            result['attachment_count'] += 1
                else:
                    if msg.get_content_type() == 'text/plain':
                        body_parts.append(msg.get_content())

                result['body'] = '\n\n'.join(body_parts)

                return result

            except Exception as e:
                logger.error(f"Failed to extract .eml file: {str(e)}")
                return result

        elif extension == '.msg':
            # Handle .msg files (Microsoft Outlook format)
            try:
                import extract_msg

                msg = extract_msg.Message(email_path)

                # Extract sender
                result['sender'] = msg.sender or ''

                # Extract recipients
                recipients = []
                if msg.to:
                    recipients.extend([addr.strip() for addr in msg.to.split(';')])
                if msg.cc:
                    recipients.extend([addr.strip() for addr in msg.cc.split(';')])
                result['recipients'] = recipients

                # Extract subject
                result['subject'] = msg.subject or ''

                # Extract date
                if msg.date:
                    result['date'] = str(msg.date)

                # Extract body
                result['body'] = msg.body or ''

                # Count attachments
                result['attachment_count'] = len(msg.attachments)

                # Check for HTML body
                result['has_html'] = bool(msg.htmlBody)

                msg.close()
                return result

            except Exception as e:
                logger.error(f"Failed to extract .msg file: {str(e)}")
                return result

        return result

    except Exception as e:
        logger.error(f"Error extracting email: {str(e)}")
        return {
            'sender': None,
            'recipients': [],
            'subject': None,
            'date': None,
            'body': '',
            'attachment_count': 0,
            'has_html': False
        }


def extract_archive_metadata(archive_path: str) -> dict:
    """
    Extract metadata from archive files (.zip, .rar, .7z, .tar, .gz).
    Does not extract contents, only analyzes the archive structure.

    Returns:
        dict with file_count, total_size, file_types, and file_list
    """
    try:
        from pathlib import Path
        import zipfile
        import tarfile

        archive_file_path = Path(archive_path)
        extension = archive_file_path.suffix.lower()

        result = {
            'file_count': 0,
            'total_size': 0,
            'file_types': {},
            'file_list': []
        }

        if extension == '.zip':
            # Handle ZIP files
            try:
                with zipfile.ZipFile(archive_path, 'r') as zf:
                    file_list = zf.namelist()
                    result['file_count'] = len(file_list)

                    # Analyze file types and sizes
                    for file_info in zf.infolist():
                        if not file_info.is_dir():
                            # Get file extension
                            file_ext = Path(file_info.filename).suffix.lower()
                            if file_ext:
                                result['file_types'][file_ext] = result['file_types'].get(file_ext, 0) + 1

                            # Add to total size
                            result['total_size'] += file_info.file_size

                            # Add to file list (limit to first 50 files)
                            if len(result['file_list']) < 50:
                                result['file_list'].append({
                                    'name': file_info.filename,
                                    'size': file_info.file_size
                                })

                return result

            except Exception as e:
                logger.error(f"Failed to extract ZIP metadata: {str(e)}")
                return result

        elif extension in ['.tar', '.gz', '.tgz']:
            # Handle TAR files (including .tar.gz)
            try:
                mode = 'r:gz' if extension in ['.gz', '.tgz'] or archive_path.endswith('.tar.gz') else 'r'
                with tarfile.open(archive_path, mode) as tf:
                    members = tf.getmembers()
                    result['file_count'] = len([m for m in members if m.isfile()])

                    # Analyze file types and sizes
                    for member in members:
                        if member.isfile():
                            # Get file extension
                            file_ext = Path(member.name).suffix.lower()
                            if file_ext:
                                result['file_types'][file_ext] = result['file_types'].get(file_ext, 0) + 1

                            # Add to total size
                            result['total_size'] += member.size

                            # Add to file list (limit to first 50 files)
                            if len(result['file_list']) < 50:
                                result['file_list'].append({
                                    'name': member.name,
                                    'size': member.size
                                })

                return result

            except Exception as e:
                logger.error(f"Failed to extract TAR metadata: {str(e)}")
                return result

        elif extension == '.7z':
            # Handle 7Z files
            try:
                import py7zr

                with py7zr.SevenZipFile(archive_path, 'r') as szf:
                    file_list = szf.getnames()
                    result['file_count'] = len(file_list)

                    # Analyze file types
                    for filename in file_list:
                        file_ext = Path(filename).suffix.lower()
                        if file_ext:
                            result['file_types'][file_ext] = result['file_types'].get(file_ext, 0) + 1

                        # Add to file list (limit to first 50 files)
                        if len(result['file_list']) < 50:
                            result['file_list'].append({
                                'name': filename,
                                'size': 0  # py7zr doesn't provide easy size access without extraction
                            })

                return result

            except Exception as e:
                logger.error(f"Failed to extract 7Z metadata: {str(e)}")
                return result

        elif extension == '.rar':
            # Handle RAR files
            try:
                import rarfile

                with rarfile.RarFile(archive_path, 'r') as rf:
                    file_list = rf.namelist()
                    result['file_count'] = len(file_list)

                    # Analyze file types and sizes
                    for file_info in rf.infolist():
                        if not file_info.isdir():
                            # Get file extension
                            file_ext = Path(file_info.filename).suffix.lower()
                            if file_ext:
                                result['file_types'][file_ext] = result['file_types'].get(file_ext, 0) + 1

                            # Add to total size
                            result['total_size'] += file_info.file_size

                            # Add to file list (limit to first 50 files)
                            if len(result['file_list']) < 50:
                                result['file_list'].append({
                                    'name': file_info.filename,
                                    'size': file_info.file_size
                                })

                return result

            except Exception as e:
                logger.error(f"Failed to extract RAR metadata: {str(e)}")
                return result

        return result

    except Exception as e:
        logger.error(f"Error extracting archive metadata: {str(e)}")
        return {
            'file_count': 0,
            'total_size': 0,
            'file_types': {},
            'file_list': []
        }


def analyze_code_file(code_path: str) -> dict:
    """
    Analyze code files to extract metadata (language, line count, etc.).

    Returns:
        dict with language, line_count, code_lines, comment_lines, blank_lines
    """
    try:
        from pathlib import Path
        from pygments.lexers import get_lexer_for_filename, guess_lexer
        from pygments.util import ClassNotFound
        import chardet

        code_file_path = Path(code_path)

        result = {
            'language': 'unknown',
            'line_count': 0,
            'code_lines': 0,
            'comment_lines': 0,
            'blank_lines': 0,
            'file_size': 0,
            'encoding': 'utf-8',
            'extracted_text': ''
        }

        # Get file size
        result['file_size'] = code_file_path.stat().st_size

        # Detect encoding
        try:
            with open(code_path, 'rb') as f:
                raw_data = f.read()
                detected = chardet.detect(raw_data)
                result['encoding'] = detected.get('encoding', 'utf-8') or 'utf-8'
        except:
            result['encoding'] = 'utf-8'

        # Read file content
        try:
            with open(code_path, 'r', encoding=result['encoding']) as f:
                content = f.read()
        except:
            # Fallback to utf-8 with errors ignored
            with open(code_path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

        # Detect language using Pygments
        try:
            lexer = get_lexer_for_filename(code_path)
            result['language'] = lexer.name
        except ClassNotFound:
            # Try to guess from content
            try:
                lexer = guess_lexer(content)
                result['language'] = lexer.name
            except:
                result['language'] = 'text'

        # Analyze lines
        lines = content.split('\n')
        result['line_count'] = len(lines)

        # Simple heuristics for code/comment/blank lines
        for line in lines:
            stripped = line.strip()
            if not stripped:
                result['blank_lines'] += 1
            elif stripped.startswith(('#', '//', '/*', '*', '<!--', '--', '%', ';')):
                result['comment_lines'] += 1
            else:
                result['code_lines'] += 1

        # Extract text content for searchability (limit to 50KB for database storage)
        result['extracted_text'] = content[:50000] if len(content) > 50000 else content

        return result

    except Exception as e:
        logger.error(f"Error analyzing code file: {str(e)}")
        return {
            'language': 'unknown',
            'line_count': 0,
            'code_lines': 0,
            'comment_lines': 0,
            'blank_lines': 0,
            'file_size': 0,
            'encoding': 'utf-8',
            'extracted_text': ''
        }


def get_paddle_ocr():
    """Lazily initialize PaddleOCR (it's heavy to load)."""
    global PADDLE_OCR
    if PADDLE_OCR is None and PADDLEOCR_AVAILABLE:
        logger.info("Initializing PaddleOCR...")
        PADDLE_OCR = PaddleOCR(use_angle_cls=True, lang='en', show_log=False)
        logger.info("PaddleOCR initialized")
    return PADDLE_OCR


def perform_ocr_with_paddle(image) -> str:
    """Perform OCR using PaddleOCR on a PIL Image."""
    ocr = get_paddle_ocr()
    if ocr is None:
        return ""

    # Convert PIL Image to numpy array for PaddleOCR
    img_array = np.array(image)

    # PaddleOCR returns list of results: [[box, (text, confidence)], ...]
    result = ocr.ocr(img_array, cls=True)

    if not result or not result[0]:
        return ""

    # Extract text from results
    lines = []
    for line in result[0]:
        if line and len(line) >= 2:
            text = line[1][0] if isinstance(line[1], tuple) else str(line[1])
            lines.append(text)

    return "\n".join(lines)


def perform_ocr_with_tesseract(image) -> str:
    """Perform OCR using Tesseract on a PIL Image."""
    if not TESSERACT_AVAILABLE:
        return ""
    return pytesseract.image_to_string(image)


def perform_ocr(document_path: str, engine: str = "auto") -> str:
    """
    Perform OCR on document (PDF or image).
    For PDFs, converts pages to images first, then performs OCR on each page.

    Args:
        document_path: Path to the document file
        engine: OCR engine to use - "auto", "paddleocr", or "tesseract"
                "auto" prefers PaddleOCR if available, falls back to Tesseract
    """
    if not OCR_AVAILABLE:
        logger.warning("No OCR engine available")
        return ""

    # Determine which OCR function to use
    if engine == "paddleocr":
        if not PADDLEOCR_AVAILABLE:
            logger.warning("PaddleOCR requested but not available, falling back to Tesseract")
            ocr_func = perform_ocr_with_tesseract
            engine_name = "Tesseract"
        else:
            ocr_func = perform_ocr_with_paddle
            engine_name = "PaddleOCR"
    elif engine == "tesseract":
        if not TESSERACT_AVAILABLE:
            logger.warning("Tesseract requested but not available, falling back to PaddleOCR")
            ocr_func = perform_ocr_with_paddle
            engine_name = "PaddleOCR"
        else:
            ocr_func = perform_ocr_with_tesseract
            engine_name = "Tesseract"
    else:  # auto - prefer PaddleOCR for better accuracy
        if PADDLEOCR_AVAILABLE:
            ocr_func = perform_ocr_with_paddle
            engine_name = "PaddleOCR"
        elif TESSERACT_AVAILABLE:
            ocr_func = perform_ocr_with_tesseract
            engine_name = "Tesseract"
        else:
            return ""

    logger.info(f"Using {engine_name} for OCR")

    try:
        doc_path = Path(document_path)
        mime_type = doc_path.suffix.lower()

        # Handle PDF files by converting to images first
        if mime_type == '.pdf':
            all_text = []

            # Try pdf2image first
            try:
                from pdf2image import convert_from_path
                logger.info(f"Converting PDF to images for OCR: {document_path}")
                images = convert_from_path(str(document_path))

                for page_num, image in enumerate(images, 1):
                    logger.info(f"Performing OCR on page {page_num}/{len(images)} with {engine_name}")
                    page_text = ocr_func(image)
                    if page_text.strip():
                        all_text.append(f"--- Page {page_num} ---\n{page_text.strip()}")

                result = "\n\n".join(all_text)
                logger.info(f"OCR completed with {engine_name}: extracted {len(result)} characters from {len(images)} pages")
                return result

            except ImportError:
                logger.warning("pdf2image not available, trying PyMuPDF for OCR")
                try:
                    import fitz  # PyMuPDF
                    doc = fitz.open(document_path)
                    logger.info(f"Converting PDF to images using PyMuPDF for OCR: {document_path}")

                    for page_num in range(len(doc)):
                        page = doc[page_num]
                        # Render page to pixmap with higher resolution for better OCR
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # 2x zoom
                        # Convert to PIL Image
                        img_data = pix.tobytes("ppm")
                        image = Image.frombytes("RGB", [pix.width, pix.height], img_data)

                        logger.info(f"Performing OCR on page {page_num + 1}/{len(doc)} with {engine_name}")
                        page_text = ocr_func(image)
                        if page_text.strip():
                            all_text.append(f"--- Page {page_num + 1} ---\n{page_text.strip()}")

                    doc.close()
                    result = "\n\n".join(all_text)
                    logger.info(f"OCR completed with {engine_name}: extracted {len(result)} characters from {len(doc)} pages")
                    return result

                except ImportError:
                    logger.error("Neither pdf2image nor PyMuPDF available for PDF OCR")
                    return ""

        # Handle regular image files
        else:
            logger.info(f"Performing OCR on image with {engine_name}: {document_path}")
            image = Image.open(document_path)
            result = ocr_func(image).strip()
            logger.info(f"OCR completed with {engine_name}: extracted {len(result)} characters")
            return result

    except Exception as e:
        logger.error(f"OCR failed: {str(e)}")
        return ""


def extract_keywords(text: str, max_keywords: int = 10) -> List[str]:
    """Extract keywords from text."""
    words = text.lower().split()
    stop_words = {'a', 'an', 'the', 'is', 'are', 'was', 'were', 'in', 'on', 'at', 'to', 'for', 'of', 'with', 'that', 'this', 'by', 'from', 'as', 'be', 'or', 'and'}
    keywords = [w.strip('.,!?;:') for w in words if w not in stop_words and len(w) > 3]
    # Count frequency and return top keywords
    from collections import Counter
    keyword_counts = Counter(keywords)
    return [word for word, count in keyword_counts.most_common(max_keywords)]


def classify_document_with_ollama(text: str, filename: str, ollama_model: str = "llama3.2") -> dict:
    """
    Classify document type using Ollama LLM.

    Args:
        text: Extracted text from document
        filename: Original filename
        ollama_model: Ollama model to use

    Returns:
        dict with document_type and confidence
    """
    global ollama_client

    if not OLLAMA_AVAILABLE or not text.strip():
        return {"document_type": "unknown", "confidence": 0.0}

    try:
        prompt = f"""Analyze this document and classify its type. Consider the filename and content.

Filename: {filename}

Content sample (first 500 chars):
{text[:500]}

Classify this document into ONE of these categories:
- invoice
- receipt
- contract
- letter
- report
- resume
- form
- certificate
- statement
- manual
- presentation
- spreadsheet
- note
- other

Respond in JSON format only: {{"document_type": "category", "confidence": 0.95}}"""

        response = ollama_client.generate(
            model=ollama_model,
            prompt=prompt,
            stream=False
        )

        result_text = response.get('response', '')

        # Parse JSON response using robust extraction
        result = extract_json_from_response(result_text)
        if result:
            return {
                "document_type": result.get("document_type", "other"),
                "confidence": float(result.get("confidence", 0.7))
            }

        return {"document_type": "other", "confidence": 0.5}

    except Exception as e:
        logger.error(f"Document classification failed: {str(e)}")
        return {"document_type": "unknown", "confidence": 0.0}


def extract_entities_with_ollama(text: str, doc_type: str, ollama_model: str = "llama3.2") -> dict:
    """
    Extract important entities from document using Ollama LLM.

    Args:
        text: Extracted text from document
        doc_type: Classified document type
        ollama_model: Ollama model to use

    Returns:
        dict with extracted entities (dates, amounts, parties, document_numbers)
    """
    global ollama_client

    if not OLLAMA_AVAILABLE or not text.strip():
        return {}

    try:
        # Type-specific extraction prompts
        if doc_type in ['invoice', 'receipt']:
            entity_types = "invoice_number, date, due_date, from_party, to_party, total_amount, currency"
        elif doc_type == 'contract':
            entity_types = "contract_number, effective_date, expiry_date, parties, value, terms"
        elif doc_type in ['letter', 'certificate']:
            entity_types = "date, from_party, to_party, reference_number, validity_date"
        else:
            entity_types = "dates, parties, reference_numbers, amounts"

        prompt = f"""Extract key information from this {doc_type} document.

Content (first 800 chars):
{text[:800]}

Extract these entities: {entity_types}

Respond in JSON format with actual values found (use null if not found):
{{
  "dates": ["2025-11-24"],
  "amounts": [{{"value": 800, "currency": "USD"}}],
  "parties": ["from_name", "to_name"],
  "document_number": "INV-008",
  "reference": "REF-123"
}}"""

        response = ollama_client.generate(
            model=ollama_model,
            prompt=prompt,
            stream=False
        )

        result_text = response.get('response', '')

        # Parse JSON response using robust extraction
        entities = extract_json_from_response(result_text)
        return entities if entities else {}

    except Exception as e:
        logger.error(f"Entity extraction failed: {str(e)}")
        return {}


def generate_intelligent_summary(text: str, doc_type: str, entities: dict, ollama_model: str = "llama3.2") -> str:
    """
    Generate intelligent summary using Ollama LLM based on document type and extracted entities.

    Args:
        text: Extracted text from document
        doc_type: Classified document type
        entities: Extracted entities
        ollama_model: Ollama model to use

    Returns:
        Concise 2-3 sentence intelligent summary
    """
    global ollama_client

    if not OLLAMA_AVAILABLE or not text.strip():
        return text[:200] + "..." if len(text) > 200 else text

    try:
        # Build entity context
        entity_context = ""
        if entities:
            entity_context = f"\n\nKey information extracted:\n{json.dumps(entities, indent=2)}"

        prompt = f"""Summarize this {doc_type} document in 2-3 concise sentences focusing on the most important information.

Document type: {doc_type}
{entity_context}

Content (first 600 chars):
{text[:600]}

Write a clear, actionable summary that highlights:
1. What this document is (type and purpose)
2. Key parties involved
3. Important dates, amounts, or deadlines
4. Any action items or critical information

Summary (2-3 sentences):"""

        response = ollama_client.generate(
            model=ollama_model,
            prompt=prompt,
            stream=False
        )

        summary = response.get('response', '').strip()

        # Fallback if summary is too short or empty
        if len(summary) < 20:
            # Create basic summary from entities
            if doc_type == 'invoice' and entities:
                doc_num = entities.get('document_number', 'N/A')
                from_party = entities.get('parties', ['Unknown'])[0] if entities.get('parties') else 'Unknown'
                to_party = entities.get('parties', ['Unknown'])[1] if entities.get('parties', [None, None])[1] else 'Unknown'
                amount = entities.get('amounts', [{}])[0] if entities.get('amounts') else {}
                amount_str = f"${amount.get('value', 0)} {amount.get('currency', 'USD')}" if amount else "N/A"
                dates = entities.get('dates', [])
                due_date = dates[0] if dates else 'N/A'
                return f"Invoice {doc_num} from {from_party} to {to_party} for {amount_str}, due {due_date}"
            else:
                return text[:200] + "..." if len(text) > 200 else text

        return summary

    except Exception as e:
        logger.error(f"Summary generation failed: {str(e)}")
        return text[:200] + "..." if len(text) > 200 else text


# ===== AUDIO PROCESSING =====

def generate_audio_thumbnail(audio_path: str, max_size: tuple = (800, 800)) -> Optional[str]:
    """
    Generate a waveform visualization thumbnail for an audio file.

    Args:
        audio_path: Path to the audio file
        max_size: Maximum dimensions for the thumbnail

    Returns:
        Path to the generated thumbnail or None if generation fails
    """
    try:
        import librosa
        import librosa.display
        import matplotlib.pyplot as plt
        import matplotlib
        matplotlib.use('Agg')  # Non-interactive backend

        # Parse the original path
        original_path = Path(audio_path)

        # Create thumbnail directory
        thumbnail_dir = original_path.parent / 'thumbnails'
        thumbnail_dir.mkdir(parents=True, exist_ok=True)

        # Generate thumbnail filename
        thumbnail_filename = original_path.stem + '.jpg'
        thumbnail_path = thumbnail_dir / thumbnail_filename

        # Load audio file (librosa automatically resamples)
        y, sr = librosa.load(audio_path, duration=60)  # Load first 60 seconds max

        # Create figure with dark theme
        fig, ax = plt.subplots(figsize=(10, 4), facecolor='#1a1a1a')
        ax.set_facecolor('#1a1a1a')

        # Plot waveform
        librosa.display.waveshow(y, sr=sr, ax=ax, color='#00d4ff', alpha=0.8)

        # Styling
        ax.set_xlabel('Time (s)', color='white', fontsize=10)
        ax.set_ylabel('Amplitude', color='white', fontsize=10)
        ax.tick_params(colors='white', labelsize=8)
        ax.grid(True, alpha=0.2, color='white')

        # Remove top and right spines
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_color('white')
        ax.spines['bottom'].set_color('white')

        # Tight layout
        plt.tight_layout()

        # Save as JPEG
        plt.savefig(str(thumbnail_path),
                   format='jpg',
                   dpi=100,
                   bbox_inches='tight',
                   facecolor='#1a1a1a',
                   edgecolor='none')
        plt.close(fig)

        logger.info(f"Generated audio waveform thumbnail: {thumbnail_path}")
        return str(thumbnail_path)

    except ImportError as e:
        logger.warning(f"Audio visualization libraries not available: {str(e)}")
        return None
    except Exception as e:
        logger.error(f"Failed to generate audio thumbnail for {audio_path}: {str(e)}")
        return None


def transcribe_audio(audio_path: str, language: Optional[str] = None) -> Dict:
    """Transcribe audio using Whisper."""
    if not WHISPER_AVAILABLE or whisper_model is None:
        return {"text": "", "language": "unknown", "confidence": 0.0}

    try:
        result = whisper_model.transcribe(
            audio_path,
            language=language,
            fp16=False
        )

        return {
            "text": result["text"].strip(),
            "language": result.get("language", "unknown"),
            "confidence": 0.9  # Whisper doesn't provide confidence, use default
        }

    except Exception as e:
        logger.error(f"Audio transcription failed: {str(e)}")
        return {"text": "", "language": "unknown", "confidence": 0.0}


# ===== TEXT EMBEDDING =====

def generate_text_embedding(text: str) -> np.ndarray:
    """Generate normalized embedding for text using CLIP."""
    # CLIP has max sequence length of 77 tokens, so truncate if needed
    inputs = clip_processor(text=[text], return_tensors="pt", padding=True, truncation=True, max_length=77).to(device)

    with torch.no_grad():
        text_features = clip_model.get_text_features(**inputs)

    embedding = text_features / text_features.norm(dim=-1, keepdim=True)
    return embedding.cpu().numpy().flatten()


# ===== API ENDPOINTS =====

@app.get("/health")
async def health_check():
    """Health check endpoint."""
    models_loaded = all([
        blip_processor is not None,
        blip_model is not None,
        clip_processor is not None,
        clip_model is not None
    ])

    return {
        "status": "healthy" if models_loaded else "initializing",
        "models_loaded": models_loaded,
        "device": str(device) if device else "unknown",
        "features": {
            "ollama": OLLAMA_AVAILABLE,
            "whisper": WHISPER_AVAILABLE and whisper_model is not None,
            "tesseract": TESSERACT_AVAILABLE
        }
    }


@app.post("/analyze-image", response_model=AnalyzeImageResponse)
async def analyze_image(request: AnalyzeImageRequest):
    """Analyze image with CLIP embeddings and face detection."""
    try:
        image_path = Path(request.image_path)
        if not image_path.exists():
            raise HTTPException(status_code=404, detail=f"Image not found: {request.image_path}")

        logger.info(f"Analyzing image: {request.image_path}")

        # Check if file is SVG (vector graphic - handle as text, not raster image)
        # IMPORTANT: SVG check MUST come BEFORE model validation since SVG files don't need models
        is_svg = (
            image_path.suffix.lower() == '.svg' or
            image_path.name.lower().endswith('.svg')
        )

        if is_svg:
            logger.info(f"Detected SVG file: {request.image_path}")
            try:
                # Extract SVG content as text (limit to 50KB for database)
                with open(image_path, 'r', encoding='utf-8', errors='ignore') as f:
                    svg_content = f.read()

                # Get file size info
                file_size = image_path.stat().st_size
                content_length = len(svg_content)

                # Truncate if too large
                extracted_text = svg_content[:50000] if content_length > 50000 else svg_content

                description = f"SVG vector image ({content_length} characters, {file_size} bytes)"

                logger.info(f"SVG processed: {content_length} chars extracted")

                return AnalyzeImageResponse(
                    description=description,
                    detailed_description=f"SVG file containing {content_length} characters of vector graphics markup",
                    meta_tags=['svg', 'vector', 'graphic', 'scalable'],
                    embedding=None,
                    faces_detected=0,
                    face_locations=[],
                    face_encodings=[],
                    thumbnail_path=None,
                    extracted_text=extracted_text
                )
            except Exception as e:
                logger.error(f"Error processing SVG file: {str(e)}")
                raise HTTPException(status_code=500, detail=f"Failed to process SVG: {str(e)}")

        # Model validation for raster images (after SVG check since SVG doesn't need models)
        if blip_model is None or clip_model is None:
            raise HTTPException(status_code=503, detail="Models not loaded yet")

        image = Image.open(image_path).convert("RGB")

        # Generate caption using selected model
        logger.info(f"Generating caption with model: {request.captioning_model}")
        caption = generate_caption(image, model=request.captioning_model)

        # Generate embedding using selected model
        logger.info(f"Generating embedding with model: {request.embedding_model}")
        embedding = generate_image_embedding(image, model=request.embedding_model)

        # Detect faces
        face_info = {"count": 0, "locations": [], "encodings": []}
        if request.detect_faces:
            face_info = detect_faces(image)

        # Generate browser-compatible thumbnail
        # This converts HEIC and other formats to JPEG for web display
        thumbnail_path = generate_thumbnail(str(image_path))

        # ============================================================
        # Maximum Analysis Coverage Features
        # ============================================================

        # Object detection via Florence-2 <OD> task (zero additional memory)
        objects_detected = None
        if request.detect_objects:
            logger.info("Running object detection with Florence-2 <OD>")
            objects_detected = detect_objects_florence(image)

        # Dominant color extraction via K-means clustering
        dominant_colors = None
        if request.extract_colors:
            logger.info("Extracting dominant colors")
            dominant_colors = extract_dominant_colors(image)

        # Image quality analysis via OpenCV
        image_quality = None
        quality_tier = None
        if request.analyze_quality:
            logger.info("Analyzing image quality")
            image_quality = analyze_image_quality(image)
            quality_tier = image_quality.get("quality_tier")

        # Perceptual hashing for duplicate detection
        phash = None
        dhash = None
        if request.compute_hashes:
            logger.info("Computing perceptual hashes")
            hashes = compute_perceptual_hashes(image)
            phash = hashes.get("phash")
            dhash = hashes.get("dhash")

        # Scene classification via Ollama (only if Ollama is enabled)
        scene_classification = None
        if request.classify_scene and request.use_ollama:
            logger.info(f"Classifying scene with Ollama ({request.ollama_model})")
            scene_classification = classify_scene_ollama(image, request.ollama_model)

        # Use Florence-2 OD labels for better semantic tags (e.g., "boat", "building")
        # instead of extract_keywords which just splits caption words ("there", "many")
        od_labels = objects_detected.get('labels', []) if objects_detected else []
        meta_tags = od_labels if od_labels else extract_keywords(caption)

        return AnalyzeImageResponse(
            description=caption,
            detailed_description=caption,
            meta_tags=meta_tags,
            embedding=embedding.tolist(),
            faces_detected=face_info["count"],
            face_locations=face_info["locations"],
            face_encodings=face_info["encodings"],
            thumbnail_path=thumbnail_path,
            # Maximum analysis coverage fields
            objects_detected=objects_detected,
            scene_classification=scene_classification,
            dominant_colors=dominant_colors,
            image_quality=image_quality,
            quality_tier=quality_tier,
            phash=phash,
            dhash=dhash
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing image: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/analyze-comprehensive", response_model=AnalyzeImageComprehensiveResponse)
async def analyze_image_comprehensive(request: AnalyzeImageComprehensiveRequest):
    """
    Comprehensive 4-pass VLM image analysis for maximum insight extraction.

    Pass 1: Content & Scene Analysis (subjects, objects, setting, composition)
    Pass 2: People & Emotion Analysis (demographics, emotions, interactions)
    Pass 3: Quality & Technical Analysis (focus, exposure, lighting, color)
    Pass 4: Context & Metadata Generation (tags, albums, keywords, occasion)

    Modes:
    - comprehensive: 4-pass analysis (~40-60 seconds per image)
    - quick: Single-pass analysis (~8 seconds per image)
    """
    try:
        # Check prerequisites
        if not OLLAMA_AVAILABLE:
            raise HTTPException(status_code=503, detail="Ollama not available - required for comprehensive analysis")

        if not COMPREHENSIVE_ANALYZER_AVAILABLE:
            raise HTTPException(status_code=503, detail="ComprehensiveImageAnalyzer not available")

        image_path = Path(request.image_path)
        if not image_path.exists():
            raise HTTPException(status_code=404, detail=f"Image not found: {request.image_path}")

        logger.info(f"Comprehensive analysis starting: {request.image_path} (mode={request.analysis_mode})")

        # Load image
        image = Image.open(image_path).convert("RGB")

        # Build image metadata for context
        image_metadata = {
            "width": image.width,
            "height": image.height,
            "path": str(image_path),
            "filename": image_path.name,
        }

        # Create analyzer with appropriate mode
        quick_mode = request.analysis_mode == "quick"
        analyzer = create_analyzer(
            ollama_client=ollama_client,
            ollama_model=request.ollama_model,
            quick_mode=quick_mode,
        )

        # Run comprehensive analysis
        result = analyzer.analyze(image, image_metadata)

        # Generate embedding if requested
        embedding = None
        if request.generate_embedding:
            logger.info(f"Generating embedding with model: {request.embedding_model}")
            embedding_array = generate_image_embedding(image, model=request.embedding_model)
            embedding = embedding_array.tolist() if embedding_array is not None else None

        # Detect faces if requested
        face_info = {"count": 0, "locations": [], "encodings": []}
        if request.detect_faces:
            face_info = detect_faces(image)

        # Generate thumbnail
        thumbnail_path = generate_thumbnail(str(image_path))

        logger.info(
            f"Comprehensive analysis completed: {result.passes_completed}/{result.passes_completed + result.passes_failed} passes "
            f"in {result.total_duration_seconds:.2f}s"
        )

        return AnalyzeImageComprehensiveResponse(
            success=result.success,
            analysis_mode=request.analysis_mode,
            content_analysis=result.content if result.content else None,
            people_analysis=result.people if result.people else None,
            quality_analysis=result.quality if result.quality else None,
            context_analysis=result.context if result.context else None,
            combined_analysis=result.combined if result.combined else None,
            passes_completed=result.passes_completed,
            passes_failed=result.passes_failed,
            analysis_duration_seconds=result.total_duration_seconds,
            errors=result.errors,
            embedding=embedding,
            faces_detected=face_info["count"],
            face_locations=face_info["locations"],
            face_encodings=face_info["encodings"],
            thumbnail_path=thumbnail_path,
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error in comprehensive analysis: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/analyze-video", response_model=AnalyzeVideoResponse)
async def analyze_video(request: AnalyzeVideoRequest):
    """Analyze video with scene detection and embeddings."""
    try:
        video_path = Path(request.video_path)
        if not video_path.exists():
            raise HTTPException(status_code=404, detail=f"Video not found: {request.video_path}")

        logger.info(f"Analyzing video: {request.video_path}")

        # Get video metadata
        metadata = get_video_metadata(str(video_path))

        # Generate browser-compatible thumbnail
        # This extracts a frame from the video and converts to JPEG for web display
        thumbnail_path = generate_video_thumbnail(str(video_path))

        # Extract and analyze frames
        scene_descriptions = []
        embeddings = []

        if request.extract_frames:
            # Use smart scene detection for better performance
            # Falls back to interval-based extraction if scene detection fails
            try:
                logger.info(f"Using smart scene detection (threshold={SCENE_THRESHOLD}, min_duration={MIN_SCENE_DURATION})")
                frames = extract_video_frames_with_scene_detection(
                    str(video_path),
                    scene_threshold=SCENE_THRESHOLD,
                    min_scene_duration_frames=MIN_SCENE_DURATION
                )

                # Fallback to interval-based extraction if no frames were extracted
                if not frames:
                    logger.warning("Scene detection returned no frames, falling back to interval-based extraction")
                    frames = extract_video_frames(str(video_path), request.frame_interval)

            except Exception as e:
                logger.error(f"Scene detection failed: {str(e)}, falling back to interval-based extraction")
                frames = extract_video_frames(str(video_path), request.frame_interval)

            scene_descriptions = analyze_video_scenes(frames)

            # Generate embeddings for key frames using parallel processing
            key_frames = frames[:5]  # Use first 5 frames
            if key_frames:
                logger.info(f"Generating embeddings for {len(key_frames)} key frames with {MAX_VIDEO_WORKERS} workers")

                def _generate_frame_embedding(frame: np.ndarray) -> np.ndarray:
                    """Helper function for parallel embedding generation."""
                    pil_image = Image.fromarray(frame)
                    return generate_image_embedding(pil_image)

                with ThreadPoolExecutor(max_workers=MAX_VIDEO_WORKERS) as executor:
                    # Submit all frames for embedding generation
                    futures = [executor.submit(_generate_frame_embedding, frame) for frame in key_frames]

                    # Collect results as they complete
                    for future in as_completed(futures):
                        try:
                            emb = future.result()
                            embeddings.append(emb)
                        except Exception as e:
                            logger.error(f"Failed to generate embedding: {str(e)}")

        # Average embeddings
        avg_embedding = np.mean(embeddings, axis=0) if embeddings else np.zeros(512)

        return AnalyzeVideoResponse(
            duration_seconds=metadata.get("duration_seconds", 0),
            frame_count=metadata.get("frame_count", 0),
            fps=metadata.get("fps", 0),
            resolution=metadata.get("resolution", "unknown"),
            scene_descriptions=scene_descriptions,
            embedding=avg_embedding.tolist(),
            objects_detected=[],
            thumbnail_path=thumbnail_path
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing video: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


def extract_pdf_text(pdf_path: str) -> str:
    """Extract native text from PDF (not OCR)."""
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(pdf_path)
        text_parts = []

        for page_num in range(len(doc)):
            page = doc[page_num]
            text = page.get_text()
            if text.strip():
                text_parts.append(text.strip())

        doc.close()
        result = "\n\n".join(text_parts)
        logger.info(f"Extracted {len(result)} characters of native text from PDF")
        return result

    except Exception as e:
        logger.error(f"PDF text extraction failed: {str(e)}")
        return ""


def extract_text_document(doc_path: str) -> str:
    """Extract text from plain text files (.txt, .md, .log, etc.)."""
    try:
        with open(doc_path, 'r', encoding='utf-8') as f:
            text = f.read()
        logger.info(f"Extracted {len(text)} characters from text document")
        return text
    except UnicodeDecodeError:
        # Try with different encoding
        try:
            with open(doc_path, 'r', encoding='latin-1') as f:
                text = f.read()
            logger.info(f"Extracted {len(text)} characters from text document (latin-1)")
            return text
        except Exception as e:
            logger.error(f"Text document extraction failed: {str(e)}")
            return ""
    except Exception as e:
        logger.error(f"Text document extraction failed: {str(e)}")
        return ""


@app.post("/analyze-document", response_model=AnalyzeDocumentResponse)
async def analyze_document(request: AnalyzeDocumentRequest):
    """Analyze document with OCR and text extraction."""
    try:
        doc_path = Path(request.document_path)
        if not doc_path.exists():
            raise HTTPException(status_code=404, detail=f"Document not found: {request.document_path}")

        logger.info(f"Analyzing document: {request.document_path}")

        # Generate browser-compatible thumbnail
        # This renders the first page of PDFs or creates a preview for text files
        thumbnail_path = generate_document_thumbnail(str(doc_path))

        # Extract text based on file type
        extracted_text = ""
        file_extension = doc_path.suffix.lower()

        # Word documents (.docx, .doc, .rtf, .odt)
        if file_extension in ['.docx', '.doc', '.rtf', '.odt']:
            logger.info(f"Extracting text from Word document: {file_extension}")
            extracted_text = extract_word_document(str(doc_path))

        # Plain text documents
        elif file_extension in ['.txt', '.md', '.log', '.csv', '.json', '.xml']:
            logger.info(f"Extracting text from plain text document: {file_extension}")
            extracted_text = extract_text_document(str(doc_path))

        # PDF documents
        elif file_extension == '.pdf':
            logger.info("Extracting text from PDF document")
            # Try native text extraction first
            extracted_text = extract_pdf_text(str(doc_path))

            # If no text extracted and OCR is requested, perform OCR
            if not extracted_text.strip() and request.perform_ocr and OCR_AVAILABLE:
                logger.info(f"No native text found in PDF, performing OCR with engine: {request.ocr_engine}")
                extracted_text = perform_ocr(str(doc_path), engine=request.ocr_engine)

        # Image documents - perform OCR if requested
        elif request.perform_ocr and OCR_AVAILABLE:
            logger.info(f"Performing OCR on image document: {file_extension} with engine: {request.ocr_engine}")
            extracted_text = perform_ocr(str(doc_path), engine=request.ocr_engine)

        else:
            logger.warning(f"Unsupported document type for text extraction: {file_extension}")

        # Extract keywords
        keywords = extract_keywords(extracted_text) if extracted_text else []

        # Generate embedding from text
        embedding = generate_text_embedding(extracted_text) if extracted_text else np.zeros(512)

        # Intelligent document analysis using Ollama
        document_type = None
        classification_confidence = None
        entities = None
        summary = None

        # Debug: Check Ollama processing conditions
        logger.info(f"DEBUG: extracted_text length: {len(extracted_text) if extracted_text else 0}")
        logger.info(f"DEBUG: request.use_ollama: {request.use_ollama}")
        logger.info(f"DEBUG: OLLAMA_AVAILABLE: {OLLAMA_AVAILABLE}")

        if extracted_text and request.use_ollama and OLLAMA_AVAILABLE:
            logger.info("Starting Ollama intelligent analysis...")
            try:
                # 1. Classify document type
                classification = classify_document_with_ollama(extracted_text, doc_path.name, request.ollama_model)
                document_type = classification.get("document_type")
                classification_confidence = classification.get("confidence")
                logger.info(f"Classified as {document_type} (confidence: {classification_confidence})")

                # 2. Extract entities based on document type
                entities = extract_entities_with_ollama(extracted_text, document_type, request.ollama_model)
                logger.info(f"Extracted entities: {entities}")

                # 3. Generate intelligent summary
                summary = generate_intelligent_summary(extracted_text, document_type, entities, request.ollama_model)
                logger.info(f"Generated summary: {summary[:100]}...")

            except Exception as e:
                logger.error(f"Intelligent document analysis failed: {str(e)}")
                # Fallback to basic summary
                summary = extracted_text[:200] + "..." if len(extracted_text) > 200 else extracted_text
        else:
            logger.warning(f"Skipping Ollama intelligent analysis - extracted_text={bool(extracted_text)}, use_ollama={request.use_ollama}, OLLAMA_AVAILABLE={OLLAMA_AVAILABLE}")

        return AnalyzeDocumentResponse(
            extracted_text=extracted_text,
            page_count=None,
            summary=summary,
            keywords=keywords,
            embedding=embedding.tolist(),
            thumbnail_path=thumbnail_path,
            document_type=document_type,
            classification_confidence=classification_confidence,
            entities=entities
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing document: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/transcribe-audio", response_model=TranscribeAudioResponse)
async def transcribe_audio_endpoint(request: TranscribeAudioRequest):
    """Transcribe audio to text using Whisper."""
    try:
        if not WHISPER_AVAILABLE:
            raise HTTPException(status_code=503, detail="Whisper not available")

        audio_path = Path(request.audio_path)
        if not audio_path.exists():
            raise HTTPException(status_code=404, detail=f"Audio not found: {request.audio_path}")

        logger.info(f"Transcribing audio: {request.audio_path}")

        # Generate waveform visualization thumbnail
        # This creates a visual representation of the audio waveform
        thumbnail_path = generate_audio_thumbnail(str(audio_path))

        # Transcribe audio
        result = transcribe_audio(str(audio_path), request.language)

        # Generate embedding from transcribed text
        embedding = generate_text_embedding(result["text"]) if result["text"] else np.zeros(512)

        return TranscribeAudioResponse(
            text=result["text"],
            language=result["language"],
            confidence=result["confidence"],
            embedding=embedding.tolist(),
            thumbnail_path=thumbnail_path
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error transcribing audio: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/embed-text", response_model=EmbedTextResponse)
async def embed_text(request: EmbedTextRequest):
    """Generate embedding for text query."""
    try:
        if clip_model is None:
            raise HTTPException(status_code=503, detail="Models not loaded yet")

        logger.info(f"Embedding text query: {request.query}")
        embedding = generate_text_embedding(request.query)

        return EmbedTextResponse(embedding=embedding.tolist())

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error embedding text: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/extract-email", response_model=ExtractEmailResponse)
async def extract_email_endpoint(request: ExtractEmailRequest):
    """Extract metadata and content from email files (.eml, .msg)."""
    try:
        email_path = Path(request.file_path)
        if not email_path.exists():
            raise HTTPException(status_code=404, detail=f"Email file not found: {request.file_path}")

        logger.info(f"Extracting email: {request.file_path}")

        result = extract_email(str(email_path))

        return ExtractEmailResponse(**result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error extracting email: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/extract-archive-metadata", response_model=ExtractArchiveMetadataResponse)
async def extract_archive_metadata_endpoint(request: ExtractArchiveMetadataRequest):
    """Extract metadata from archive files (.zip, .rar, .7z, .tar, .gz)."""
    try:
        archive_path = Path(request.file_path)
        if not archive_path.exists():
            raise HTTPException(status_code=404, detail=f"Archive file not found: {request.file_path}")

        logger.info(f"Extracting archive metadata: {request.file_path}")

        result = extract_archive_metadata(str(archive_path))

        return ExtractArchiveMetadataResponse(**result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error extracting archive metadata: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.post("/analyze-code-file", response_model=AnalyzeCodeFileResponse)
async def analyze_code_file_endpoint(request: AnalyzeCodeFileRequest):
    """Analyze code files to extract metadata (language, line count, etc.)."""
    try:
        code_path = Path(request.file_path)
        if not code_path.exists():
            raise HTTPException(status_code=404, detail=f"Code file not found: {request.file_path}")

        logger.info(f"Analyzing code file: {request.file_path}")

        result = analyze_code_file(str(code_path))

        return AnalyzeCodeFileResponse(**result)

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error analyzing code file: {str(e)}")
        raise HTTPException(status_code=500, detail=str(e))


@app.get("/")
async def root():
    """Root endpoint."""
    return {
        "service": "Avinash-EYE Multi-Media AI Service",
        "version": "3.0.0",
        "features": [
            "Image captioning (BLIP)",
            "Multi-modal embeddings (CLIP)",
            "Face detection and recognition",
            "Video scene analysis",
            "Document OCR" if TESSERACT_AVAILABLE else "Document OCR (disabled)",
            "Audio transcription (Whisper)" if WHISPER_AVAILABLE else "Audio transcription (disabled)",
        ],
        "endpoints": [
            "/health",
            "/analyze-image",
            "/analyze-video",
            "/analyze-document",
            "/transcribe-audio",
            "/embed-text",
            "/extract-email",
            "/extract-archive-metadata",
            "/analyze-code-file"
        ]
    }
